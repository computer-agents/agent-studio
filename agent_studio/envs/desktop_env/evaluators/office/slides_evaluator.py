import logging
import xml.etree.ElementTree as ET

from pptx import Presentation

from agent_studio.config import Config
from agent_studio.envs.desktop_env.evaluators.evaluator import (
    Evaluator,
    FeedbackException,
    evaluation_handler,
)

config = Config()
logger = logging.getLogger(__name__)


class SlidesEvaluator(Evaluator):
    name: str = "slides"

    @evaluation_handler("compare_pptx_files")
    def compare_pptx_files(self, expected, result, options):
        # todo: not strictly match since not all information is compared because we cannot get the info through pptx  # noqa: E501
        try:
            prs1 = Presentation(expected)
        except Exception as e:
            raise FeedbackException(f"Fail to open the expected file {expected}: {e}")
        try:
            prs2 = Presentation(result)
        except Exception as e:
            raise FeedbackException(f"Fail to open the result file {result}: {e}")

        examine_number_of_slides = options.get("examine_number_of_slides", True)
        examine_shape = options.get("examine_shape", True)
        examine_text = options.get("examine_text", True)
        examine_indent = options.get("examine_indent", True)
        examine_font_name = options.get("examine_font_name", True)
        examine_font_size = options.get("examine_font_size", True)
        examine_font_bold = options.get("examine_font_bold", True)
        examine_font_italic = options.get("examine_font_italic", True)
        examine_color_rgb = options.get("examine_color_rgb", True)
        examine_font_underline = options.get("examine_font_underline", True)
        examine_strike_through = options.get("examine_strike_through", True)
        examine_alignment = options.get("examine_alignment", True)
        examine_title_bottom_position = options.get(
            "examine_title_bottom_position", False
        )
        examine_table_bottom_position = options.get(
            "examine_table_bottom_position", False
        )
        examine_right_position = options.get("examine_right_position", False)
        examine_top_position = options.get("examine_top_position", False)
        examine_shape_for_shift_size = options.get(
            "examine_shape_for_shift_size", False
        )
        examine_image_size = options.get("examine_image_size", False)
        examine_modify_height = options.get("examine_modify_height", False)
        examine_bullets = options.get("examine_bullets", True)
        examine_background_color = options.get("examine_background_color", True)
        examine_note = options.get("examine_note", True)

        # compare the number of slides
        if len(prs1.slides) != len(prs2.slides) and examine_number_of_slides:
            raise FeedbackException(
                f"Number of slides are different: {len(prs1.slides)} != {len(prs2.slides)}"  # noqa: E501
            )

        slide_idx = 0
        # compare the content of each slide
        for slide1, slide2 in zip(prs1.slides, prs2.slides):
            slide_idx += 1

            def get_slide_background_color(slide):
                background = slide.background
                if background.fill.background():
                    return background.fill.fore_color.rgb
                else:
                    return None

            if (
                get_slide_background_color(slide1) != get_slide_background_color(slide2)
                and examine_background_color
            ):
                raise FeedbackException("Background color is different")

            def get_slide_notes(slide):
                notes_slide = slide.notes_slide
                if notes_slide:
                    return notes_slide.notes_text_frame.text
                else:
                    return None

            notes1 = get_slide_notes(slide1).strip()
            notes2 = get_slide_notes(slide2).strip()
            if notes1 != notes2 and examine_note:
                raise FeedbackException(f"Notes are different: {notes1} != {notes2}")
            # check if the shapes are the same
            for shape1, shape2 in zip(slide1.shapes, slide2.shapes):
                if examine_title_bottom_position:
                    if (
                        hasattr(shape1, "text")
                        and hasattr(shape2, "text")
                        and shape1.text == shape2.text
                    ):
                        if shape1.text == "Product Comparison" and (
                            shape1.top <= shape2.top or shape1.top < 3600000
                        ):
                            raise FeedbackException(
                                "Title bottom position is different"
                            )
                    elif (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        raise FeedbackException("Title bottom position is different")

                if examine_table_bottom_position:
                    if (
                        slide_idx == 3
                        and shape1.shape_type == 19
                        and shape2.shape_type == 19
                    ):
                        if shape1.top > shape2.top or shape1.top < 3600000:
                            raise FeedbackException(
                                f"Table bottom position is different, {shape1.top} > {shape2.top} or {shape1.top} < 3600000"  # noqa: E501
                            )
                    elif (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        raise FeedbackException(
                            f"Table bottom position is different, ({shape1.left}, {shape1.top}, {shape1.width}, {shape1.height}) != ({shape2.left}, {shape2.top}, {shape2.width}, {shape2.height})"  # noqa: E501
                        )

                if examine_right_position:
                    if (
                        slide_idx == 2
                        and not hasattr(shape1, "text")
                        and not hasattr(shape2, "text")
                    ):
                        if shape1.left <= shape2.left or shape1.left < 4320000:
                            raise FeedbackException("Right position is different")

                if examine_top_position:
                    if (
                        slide_idx == 2
                        and shape1.shape_type == 13
                        and shape2.shape_type == 13
                    ):
                        if shape1.top < shape2.top or shape1.top > 1980000:
                            raise FeedbackException(
                                f"Top position is different, {shape1.top} < {shape2.top} or {shape1.top} > 1980000"  # noqa: E501
                            )
                    elif (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        raise FeedbackException(
                            f"Top position is different, {(shape1.left, shape1.top, shape1.width, shape1.height)} != {(shape2.left, shape2.top, shape2.width, shape2.height)}"  # noqa: E501
                        )

                if examine_shape_for_shift_size:
                    if (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        if not (
                            hasattr(shape1, "text")
                            and hasattr(shape2, "text")
                            and shape1.text == shape2.text
                            and shape1.text == "Elaborate on what you want to discuss."
                        ):
                            raise FeedbackException("Shape for shift size is different")

                if (
                    shape1.left != shape2.left
                    or shape1.top != shape2.top
                    or shape1.width != shape2.width
                    or shape1.height != shape2.height
                ) and examine_shape:
                    raise FeedbackException(
                        f"Shape is different: {(shape1.left, shape1.top, shape1.width, shape1.height)} != {(shape2.left, shape2.top, shape2.width, shape2.height)}"  # noqa: E501
                    )

                if examine_image_size:
                    if shape1.shape_type == 13 and shape2.shape_type == 13:
                        if (
                            shape1.width != shape2.width
                            or shape1.height != shape2.height
                        ):
                            raise FeedbackException(
                                f"Image size is different: {(shape1.width, shape1.height)} != {(shape2.width, shape2.height)}"  # noqa: E501
                            )
                    elif (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        raise FeedbackException(
                            f"Image size is different: {(shape1.left, shape1.top, shape1.width, shape1.height)} != {(shape2.left, shape2.top, shape2.width, shape2.height)}"  # noqa: E501
                        )

                if examine_modify_height:
                    if (
                        not hasattr(shape1, "text")
                        and not hasattr(shape2, "text")
                        or shape1.shape_type == 5
                        and shape2.shape_type == 5
                    ):
                        if shape1.height != shape2.height:
                            raise FeedbackException("Modify height is different")
                    elif (
                        shape1.left != shape2.left
                        or shape1.top != shape2.top
                        or shape1.width != shape2.width
                        or shape1.height != shape2.height
                    ):
                        raise FeedbackException("Modify height is different")

                if hasattr(shape1, "text") and hasattr(shape2, "text"):
                    if shape1.text.strip() != shape2.text.strip() and examine_text:
                        raise FeedbackException("Text is different")

                        # check if the paragraphs are the same
                    for para1, para2 in zip(
                        shape1.text_frame.paragraphs, shape2.text_frame.paragraphs
                    ):
                        if para1.alignment != para2.alignment and examine_alignment:
                            raise FeedbackException(
                                f"Alignment is different: {para1.alignment} != {para2.alignment}"  # noqa: E501
                            )

                        # check if the runs are the same
                        if para1.text != para2.text and examine_text:
                            raise FeedbackException("Text is different")

                        if para1.level != para2.level and examine_indent:
                            raise FeedbackException(
                                f"Indent is different, {para1.level} != {para2.level}"
                            )

                        for run1, run2 in zip(para1.runs, para2.runs):
                            # check if the font properties are the same
                            if run1.font.name != run2.font.name and examine_font_name:
                                raise FeedbackException(
                                    f"Font name is different, {run1.font.name} != {run2.font.name}"  # noqa: E501
                                )

                            if run1.font.size != run2.font.size and examine_font_size:
                                raise FeedbackException(
                                    f"Font size is different, {run1.font.size} != {run2.font.size}"  # noqa: E501
                                )

                            if run1.font.bold != run2.font.bold and examine_font_bold:
                                raise FeedbackException(
                                    f"Font bold is different, {run1.font.bold} != {run2.font.bold}"  # noqa: E501
                                )

                            if (
                                run1.font.italic != run2.font.italic
                                and examine_font_italic
                            ):
                                raise FeedbackException(
                                    f"Font italic is different, {run1.font.italic} != {run2.font.italic}"  # noqa: E501
                                )

                            if hasattr(run1.font.color, "rgb") and hasattr(
                                run2.font.color, "rgb"
                            ):
                                if (
                                    run1.font.color.rgb != run2.font.color.rgb
                                    and examine_color_rgb
                                ):
                                    raise FeedbackException(
                                        f"Color is different, {run1.font.color.rgb} != {run2.font.color.rgb}"  # noqa: E501
                                    )

                            if (
                                run1.font.underline != run2.font.underline
                                and examine_font_underline
                            ):
                                raise FeedbackException(
                                    f"Underline is different, {run1.font.underline} != {run2.font.underline}"  # noqa: E501
                                )

                            if (
                                run1.font._element.attrib.get("strike", "noStrike")
                                != run2.font._element.attrib.get("strike", "noStrike")
                                and examine_strike_through
                            ):
                                raise FeedbackException(
                                    f"Strike through is different, {run1.font._element.attrib.get('strike', 'noStrike')} != {run2.font._element.attrib.get('strike', 'noStrike')}"  # noqa: E501
                                )

                            def _extract_bullets(xml_data):
                                root = ET.fromstring(xml_data)

                                namespaces = {
                                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",  # noqa: E501
                                    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",  # noqa: E501
                                }

                                bullets = []

                                for paragraph in root.findall(".//a:p", namespaces):
                                    pPr = paragraph.find("a:pPr", namespaces)
                                    if pPr is not None:
                                        lvl = pPr.get("lvl")
                                        buChar = pPr.find("a:buChar", namespaces)
                                        char = (
                                            buChar.get("char")
                                            if buChar is not None
                                            else "No Bullet"
                                        )
                                        buClr = pPr.find(
                                            "a:buClr/a:srgbClr", namespaces
                                        )
                                        color = (
                                            buClr.get("val")
                                            if buClr is not None
                                            else "No Color"
                                        )
                                    else:
                                        lvl = "No Level"
                                        char = "No Bullet"
                                        color = "No Color"

                                    text = "".join(
                                        t.text
                                        for t in paragraph.findall(".//a:t", namespaces)
                                    )
                                    if text.strip() != "":
                                        bullets.append((lvl, char, text, color))

                                return bullets

                            if examine_bullets and _extract_bullets(
                                run1.part.blob.decode("utf-8")
                            ) != _extract_bullets(run2.part.blob.decode("utf-8")):
                                raise FeedbackException(
                                    f"Bullets are different, {_extract_bullets(run1.part.blob.decode('utf-8'))} != {_extract_bullets(run2.part.blob.decode('utf-8'))}"  # noqa: E501
                                )

                        # fixme: Actually there are more properties to be compared, we can add them later via parsing the xml data  # noqa: E501
